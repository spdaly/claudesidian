#!/usr/bin/env python3
"""
ABOUTME: Creates Seven Corners AI Virtual Assistant pitch deck PowerPoint
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define layout indices
TITLE_SLIDE = 0
TITLE_AND_CONTENT = 1
BLANK = 6

def add_title_slide():
    """Slide 1: Title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_SLIDE])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Transforming Seven Corners Customer Experience with AI"
    subtitle.text = "A Phased Approach to 24/7 Intelligent Service\n\n[Your Company Name]\n[Presenter Names]\nOctober 13, 2025"

def add_moment_of_truth():
    """Slide 2: The Moment of Truth"""
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    title = slide.shapes.title
    title.text = "Your Customers Need You When It Matters Most"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Travel insurance isn't a commodity - it's protection during vulnerable moments"

    p = tf.add_paragraph()
    p.text = "Your customers span 12+ time zones, speak 50+ languages"
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Questions range from simple ("What\'s covered?") to complex ("I\'m hospitalized in Thailand - what now?")'
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "[Visual: Split screen showing stressed traveler at 2 AM + confused family reviewing policy]"
    p.level = 0
    p.font.italic = True

def add_contact_center_reality():
    """Slide 3: The Contact Center Reality"""
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    title = slide.shapes.title
    title.text = "Three Forces Reshaping Customer Service"

    content = slide.placeholders[1]
    tf = content.text_frame

    tf.text = "1. Customer Expectations"
    p = tf.add_paragraph()
    p.text = '"Instant answers, any channel, any time"'
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "2. Operational Costs"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Agent costs rising, volumes increasing"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Average cost per contact: $12 and climbing"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "3. Competitive Pressure"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Digital-first insurers setting new standards"
    p.level = 1

def add_customer_journey():
    """Slide 4: Reimagining the Customer Journey"""
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    title = slide.shapes.title
    title.text = "Reimagining the Customer Journey"

    content = slide.placeholders[1]
    tf = content.text_frame

    tf.text = "BEFORE (Current State):"
    p = tf.add_paragraph()
    p.text = "Customer searches website FAQ → Can't find answer (5 min)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Waits on hold (15 min average)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Explains situation to agent (5 min)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Gets transferred, explains AGAIN (10 min)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "TOTAL TIME: 45+ minutes | Customer Feeling: Frustrated"
    p.level = 1
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "AFTER (AI-Powered):"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Opens chat → Authenticated automatically → AI answers with context"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "TOTAL TIME: 30 seconds | Customer Feeling: Impressed"
    p.level = 1
    p.font.bold = True

def add_technology_foundation():
    """Slide 5: The Technology Foundation"""
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    title = slide.shapes.title
    title.text = "Microsoft AI Stack for Insurance"

    content = slide.placeholders[1]
    tf = content.text_frame

    tf.text = "CUSTOMER (Web, Mobile, Phone, WhatsApp)"
    p = tf.add_paragraph()
    p.text = "↓"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "COPILOT STUDIO - Conversational UI Layer"
    p.level = 0
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "↓"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "AZURE AI SEARCH - Knowledge Base (Policies, FAQs, Coverage Docs)"
    p.level = 0
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "↓"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "AZURE AI FOUNDRY - Custom AI Models (Phase 3)"
    p.level = 0
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "↓"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "YOUR BACKEND SYSTEMS (Policy Admin, CRM, Claims)"
    p.level = 0
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Why Microsoft?"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "✓ Enterprise security & compliance (HIPAA, GDPR, SOC2)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "✓ 1,000+ pre-built connectors"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "✓ Low-code + Custom AI capabilities"
    p.level = 1

def add_conversation_example():
    """Slide 6: AI That Understands Insurance"""
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    title = slide.shapes.title
    title.text = "AI That Understands Insurance Conversations"

    content = slide.placeholders[1]
    tf = content.text_frame

    tf.text = 'CUSTOMER: "I\'m supposed to fly to Paris tomorrow but I tested positive for COVID. Am I covered?"'

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "AI RESPONSE:"
    p.level = 0
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = '"I\'m sorry to hear you\'re not feeling well. Let me help you with that."'
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "I've pulled up your Trip Protection Plus policy. Yes, you're covered for trip cancellation due to illness, including COVID-19."
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Here's what happens next:"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "1. You'll need a doctor's note"
    p.level = 2

    p = tf.add_paragraph()
    p.text = "2. File a claim within 30 days"
    p.level = 2

    p = tf.add_paragraph()
    p.text = "3. Keep all receipts for non-refundable expenses"
    p.level = 2

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Why this works: Empathy + Accuracy + Action + Choice"
    p.level = 0
    p.font.bold = True

def add_phased_methodology():
    """Slide 7: Our Phased Methodology"""
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    title = slide.shapes.title
    title.text = "Our Phased Methodology: Fast Value, Controlled Risk"

    content = slide.placeholders[1]
    tf = content.text_frame

    tf.text = "PHASE 1: FOUNDATION (Months 1-4)"
    tf.paragraphs[0].font.bold = True

    p = tf.add_paragraph()
    p.text = "🎯 Goal: Quick wins, prove value"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "📦 Deliverable: Website chat handling top 10 inquiries"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "📊 Target: 30% containment rate, 20% CSAT improvement"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "💰 Investment: $150,000"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "PHASE 2: EXPANSION (Months 5-9)"
    p.level = 0
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "🎯 Goal: Scale across channels"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "📦 Deliverable: Mobile, WhatsApp, Teams + top 30 inquiries"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "💰 Investment: $175,000"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "PHASE 3: ADVANCED AI (Months 10-18)"
    p.level = 0
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "🎯 Goal: Custom intelligence, voice channel"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "📦 Deliverable: Voice capability, custom models, recommendations"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "💰 Investment: $125,000"
    p.level = 1

def add_discovery_process():
    """Slide 8: Phase 1 Deep Dive"""
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    title = slide.shapes.title
    title.text = "We Start By Learning Your Business"

    content = slide.placeholders[1]
    tf = content.text_frame

    tf.text = "MONTH 1: DISCOVERY (4 Weeks)"
    tf.paragraphs[0].font.bold = True

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Week 1-2: Data Analysis"
    p.level = 0
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "✓ Review contact center data (3-6 months)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "✓ Identify top inquiry categories and volumes"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "✓ Analyze conversation patterns"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Week 3: Collaborative Workshop"
    p.level = 0
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "✓ Full-day session with your team"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "✓ Validate assumptions and priorities"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "✓ Define success metrics together"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Week 4: Blueprint Delivery"
    p.level = 0
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "✓ Conversation design for top 10 scenarios"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "✓ Integration architecture (specific to YOUR systems)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "✓ Phase 1 detailed plan + Go/no-go decision point"
    p.level = 1

def add_differentiation():
    """Slide 9: What Makes Us Different"""
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    title = slide.shapes.title
    title.text = "The Intersection of Three Expertises"

    content = slide.placeholders[1]
    tf = content.text_frame

    tf.text = "[Visual: Venn diagram - Microsoft AI ∩ Insurance ∩ Contact Center = US]"
    tf.paragraphs[0].font.italic = True

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "1. Microsoft AI Platform Expertise"
    p.level = 0
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Certified partners, 50+ Azure AI implementations"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "2. Insurance Domain Knowledge"
    p.level = 0
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Policy language, claims processes, compliance requirements"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "3. Contact Center Transformation"
    p.level = 0
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Agent training, KPIs, hybrid human-AI operating models"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = '"Others can build you a chatbot. We build you a contact center transformation that happens to use AI."'
    p.level = 0
    p.font.bold = True
    p.font.italic = True

def add_investment():
    """Slide 10: The Investment"""
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    title = slide.shapes.title
    title.text = "Total 18-Month Investment: $450,000"

    content = slide.placeholders[1]
    tf = content.text_frame

    tf.text = "Phase 1 - FOUNDATION (4 months): $150,000"
    p = tf.add_paragraph()
    p.text = "MVP virtual assistant, top 10 inquiries, analytics dashboard"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Payment: 50% upfront, 50% at MVP launch"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Phase 2 - EXPANSION (5 months): $175,000"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Multi-channel deployment, top 30 inquiries, CRM integration"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Phase 3 - ADVANCED AI (9 months): $125,000"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Voice channel, custom AI models, personalization engine"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Risk Mitigation:"
    p.level = 0
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "✓ Phase 1 is your pilot - $150K to prove the concept"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "✓ Phases 2-3 contingent on Phase 1 success"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "✓ Go/no-go decision points between each phase"
    p.level = 1

def add_roi():
    """Slide 11: The Return"""
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    title = slide.shapes.title
    title.text = "Projected 3-Year ROI: 356%"

    content = slide.placeholders[1]
    tf = content.text_frame

    tf.text = "Year 1 Benefits: $400,000"
    tf.paragraphs[0].font.bold = True

    p = tf.add_paragraph()
    p.text = "Call volume reduction (30% of routine inquiries automated)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Agent productivity improvement (20%)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Year 2 Benefits: $850,000"
    p.level = 0
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Expanded containment (40% automated resolution)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Customer retention improvement (2% churn reduction)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Year 3 Benefits: $800,000"
    p.level = 0
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Full optimization (50% containment rate)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Voice channel automation, 30% agent productivity gain"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "3-Year Total: $2.05M benefits - $450K investment = $1.6M net"
    p.level = 0
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Payback Period: ~9 months"
    p.level = 0
    p.font.bold = True

def add_metrics():
    """Slide 12: Success Metrics Dashboard"""
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    title = slide.shapes.title
    title.text = "How We'll Measure Success Together"

    content = slide.placeholders[1]
    tf = content.text_frame

    tf.text = "Customer Experience:"
    tf.paragraphs[0].font.bold = True

    p = tf.add_paragraph()
    p.text = "📈 CSAT Score: Baseline → +25% improvement"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "⚡ Average Response Time: 50% reduction"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "🎯 First Contact Resolution: 60%+ automated"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Operational Efficiency:"
    p.level = 0
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "🤖 Containment Rate: 40%"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "⏱️ Agent Handle Time: 30% reduction"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "🌐 24/7 Availability: <2 second response time"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Business Impact:"
    p.level = 0
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "💰 Cost per Interaction: 35% reduction"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "📞 Call Volume: 30-40% reduction (routine inquiries)"
    p.level = 1

def add_why_now():
    """Slide 13: Why Act Now"""
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    title = slide.shapes.title
    title.text = "The Window Is Closing"

    content = slide.placeholders[1]
    tf = content.text_frame

    tf.text = "2024-2025: Early Adopter Advantage ← YOU ARE HERE"
    tf.paragraphs[0].font.bold = True

    p = tf.add_paragraph()
    p.text = "Technology mature, 18-24 month head start possible"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "2026: Mainstream Adoption"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Industry standard, no longer differentiator"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "2027+: Table Stakes"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Required to compete, playing catch-up vs. leading"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Four Forces Creating Urgency:"
    p.level = 0
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "🏃 Competitive Pressure - Digital-first insurers already deploying AI"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "📱 Customer Expectations - 2025 travelers expect instant service"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "💡 Technology Maturity - Production-ready Microsoft stack NOW"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "💵 Cost Efficiency - Early adopters gain learning curve advantage"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = '"In 24 months, AI-powered service won\'t be a differentiator - it will be table stakes. The question is whether Seven Corners leads or follows."'
    p.level = 0
    p.font.italic = True

def add_next_steps():
    """Slide 14: Next Steps"""
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    title = slide.shapes.title
    title.text = "Let's Validate This Together"

    content = slide.placeholders[1]
    tf = content.text_frame

    tf.text = "STEP 1: Discovery Engagement (2 weeks)"
    tf.paragraphs[0].font.bold = True

    p = tf.add_paragraph()
    p.text = "What: Review your data, interview your team, analyze systems"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Deliverable: Custom implementation blueprint FOR Seven Corners"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Investment: $15,000 (credited toward Phase 1 if you proceed)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "STEP 2: Decision Point (1 week)"
    p.level = 0
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "You review the actual plan based on YOUR data"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "You decide: Does this align with our strategy?"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "No pressure, no obligation"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "STEP 3: Phase 1 Kickoff (if aligned)"
    p.level = 0
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "4-month MVP implementation"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Defined success metrics, go/no-go before Phase 2"
    p.level = 1

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "CALL TO ACTION:"
    p.level = 0
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "📞 Let's schedule the discovery engagement"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "[Your name, email, phone]"
    p.level = 1

# Create all slides
add_title_slide()
add_moment_of_truth()
add_contact_center_reality()
add_customer_journey()
add_technology_foundation()
add_conversation_example()
add_phased_methodology()
add_discovery_process()
add_differentiation()
add_investment()
add_roi()
add_metrics()
add_why_now()
add_next_steps()

# Save the presentation
output_path = "/Users/spdaly/claudesidian/01_Projects/Seven Corners Travel Insurance - Customer Service Virtual Assistant/Seven Corners AI Virtual Assistant Pitch.pptx"
prs.save(output_path)

print(f"PowerPoint presentation created successfully: {output_path}")
print(f"Total slides: {len(prs.slides)}")
print("\nNext steps:")
print("1. Open the PowerPoint file")
print("2. Apply your company's design template/theme")
print("3. Add visuals (diagrams, icons, charts) as noted in slides")
print("4. Customize with your branding (logo, colors, fonts)")
print("5. Review and refine talking points")
